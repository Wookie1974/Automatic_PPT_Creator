import os
import json
from pptx import Presentation
from pptx.util import Inches
from urllib.parse import urljoin
import asyncio
from playwright.async_api import async_playwright

JS_GET_MAIN_AND_IMAGES = """
() => {
  function score(n){ const ps=n.querySelectorAll('p,li'); return ps.length*10 + (n.innerText||'').length; }
  const cands = Array.from(document.querySelectorAll('article,main,section,div'));
  let best=null, sc=-1;
  for (const n of cands){ const s=score(n); if(s>sc){sc=s; best=n;} }
  const container = best || document.body;

  // Assign stable indices to images in DOM order
  const imgEls = Array.from(container.querySelectorAll('img, figure img'));
  imgEls.forEach((im, i) => im.setAttribute('data-x-img-idx', String(i+1)));

  function pick(el){
    const attrs=['src','data-src','data-original','data-lazy','data-image','data-thumb'];
    for(const a of attrs){
      const v = el.getAttribute(a);
      if (v) return { url:v, alt:(el.getAttribute('alt')||'').trim() };
    }
    return null;
  }

  const imgs = [];
  imgEls.forEach((im, i) => {
    const got = pick(im);
    if (got) imgs.push({ idx: i+1, ...got });
  });

  // Title
  let title='';
  const h = container.querySelector('h1, h2') || document.querySelector('h1, h2');
  if (h) title=(h.innerText||h.textContent||'').trim();
  if (!title) title=document.title||'';

  return { title, html: container.innerHTML, images: imgs };
}
"""

async def render_and_extract(ctx, url: str) -> dict:
    page = await ctx.new_page()
    await page.goto(url, wait_until="domcontentloaded", timeout=30000)
    try: await page.wait_for_load_state("networkidle", timeout=8000)
    except: pass
    data = await page.evaluate(JS_GET_MAIN_AND_IMAGES)
    # absolutize image URLs
    for im in data["images"]:
        im["abs"] = urljoin(url, im.get("url"))
    await page.close()
    return data

def add_slide(prs, title, text, images):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    slide.shapes.title.text = title
    slide.placeholders[1].text = text[:1000]  # Limit text for demo

    # Add all images to the slide
    left = Inches(5)
    top = Inches(1.5)
    img_width = Inches(2)
    img_gap = Inches(0.2)
    for idx, img_url in enumerate(images):
        try:
            import requests
            img_data = requests.get(img_url, timeout=10).content
            img_path = f"temp_img_{idx}.jpg"
            with open(img_path, "wb") as f:
                f.write(img_data)
            slide.shapes.add_picture(img_path, left, top + idx * (img_width + img_gap), width=img_width)
            os.remove(img_path)
        except Exception as e:
            print(f"Could not add image from {img_url}: {e}")

async def main():
    # Default Chrome bookmarks path (Mac)
    bookmarks_file = os.path.expanduser("~/Library/Application Support/Google/Chrome/Default/Bookmarks")
    folder_name = input("Enter the Chrome bookmarks folder name: ").strip()

    def find_folder(node, folder_name):
        if node.get("type") == "folder" and node.get("name") == folder_name:
            return node
        for child in node.get("children", []):
            result = find_folder(child, folder_name)
            if result:
                return result
        return None

    def extract_urls(folder_node):
        urls = []
        for child in folder_node.get("children", []):
            if child.get("type") == "url":
                urls.append(child["url"])
            elif child.get("type") == "folder":
                urls.extend(extract_urls(child))
        return urls

    with open(bookmarks_file, "r", encoding="utf-8") as f:
        data = json.load(f)
    roots = data["roots"]
    urls = []
    for root in roots.values():
        folder = find_folder(root, folder_name)
        if folder:
            urls = extract_urls(folder)
            break

    prs = Presentation()
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        ctx = await browser.new_context(viewport={"width":1400,"height":1600})
        for url in urls:
            print(f"Processing: {url}")
            data = await render_and_extract(ctx, url)
            title = data.get("title") or url
            # Extract text from HTML
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(data.get("html") or "", "lxml")
            main = soup.find("article") or soup.find("main") or soup.body
            text = main.get_text(separator="\n", strip=True) if main else ""
            images = [im["abs"] for im in data.get("images", []) if im.get("abs")]
            add_slide(prs, title, text, images)
        await ctx.close(); await browser.close()
    prs.save("bookmarks_output.pptx")
    print("PowerPoint file created: bookmarks_output.pptx")

if __name__ == "__main__":
    asyncio.run(main())