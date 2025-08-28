# slidespec_to_ppt.py
# Build a PPT deck from slidespec.json using your template's placeholders.
# - Title -> TITLE placeholder (or slide title)
# - Bullets -> BODY placeholder (or any text-capable placeholder)
# - Images -> picture placeholders first; leftover images placed in a right column fallback
#
# Usage:
#   python slidespec_to_ppt.py \
#     --index out_2026_clean/slidespec_index.csv \
#     --template "/path/to/Template_for_Autobuild.pptx" \
#     --out vault_2026.pptx

import argparse, csv, json, tempfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt
from PIL import Image as PILImage

# Optional SVG support
try:
    import cairosvg
    HAS_CAIROSVG = True
except Exception:
    HAS_CAIROSVG = False

SUPPORTED_RASTERS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".emf"}

# ---------------- Helpers ----------------

def iter_slides(spec_path: Path):
    data = json.loads(spec_path.read_text(encoding="utf-8"))
    for slide in data.get("slides", []):
        yield slide

def find_preferred_layout(prs: Presentation):
    """Pick a layout that has TITLE and BODY; otherwise first layout."""
    for layout in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in layout.placeholders}
        if PP_PLACEHOLDER.TITLE in types and PP_PLACEHOLDER.BODY in types:
            return layout
    return prs.slide_layouts[0]

def get_placeholder(slide, phtype):
    """Find first placeholder on slide of a given placeholder type."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == phtype:
                return ph
        except Exception:
            continue
    return None

def ensure_supported_image(src_path: Path, tmpdir: Path) -> Path:
    ext = src_path.suffix.lower()
    if ext in SUPPORTED_RASTERS and src_path.exists():
        return src_path
    out_path = tmpdir / (src_path.stem + ".png")
    if ext == ".svg":
        if not HAS_CAIROSVG:
            raise RuntimeError(f"SVG image found but cairosvg not installed: {src_path}")
        cairosvg.svg2png(url=str(src_path), write_to=str(out_path))
        return out_path
    # generic convert via Pillow
    with PILImage.open(src_path) as im:
        im.save(out_path, format="PNG")
    return out_path

# ---------------- Title / Bullets ----------------

def apply_title(slide, text: str):
    ph_title = get_placeholder(slide, PP_PLACEHOLDER.TITLE) or getattr(slide.shapes, "title", None)
    if ph_title is None:
        # last resort: add a textbox near top
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(9), Inches(0.9))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text or ""
        return
    tf = ph_title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text or ""

def apply_bullets(slide, bullets):
    # 1) Prefer BODY placeholder
    ph_body = get_placeholder(slide, PP_PLACEHOLDER.BODY)

    # 2) Fallback: any non-title placeholder that accepts text
    if ph_body is None:
        for ph in slide.placeholders:
            try:
                pht = ph.placeholder_format.type
            except Exception:
                pht = None
            if pht in (PP_PLACEHOLDER.TITLE, getattr(PP_PLACEHOLDER, "CENTER_TITLE", None)):
                continue
            if getattr(ph, "has_text_frame", False):
                ph_body = ph
                break

    # 3) Final fallback: create a textbox in a sane left-column area
    if ph_body is None:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.8), Inches(5.5))
        tf = box.text_frame
    else:
        tf = ph_body.text_frame

    tf.clear()
    if not bullets:
        return
    for i, b in enumerate(bullets or []):
        para = tf.add_paragraph() if i else tf.paragraphs[0]
        para.text = b or ""
        para.level = 0

# ---------------- Image Placement ----------------

def find_image_placeholders(slide):
    """
    Return a prioritized list of placeholders for images:
      1) True PICTURE placeholders
      2) Placeholders whose name includes 'picture' or 'image'
      3) Any non-title, non-body placeholder
    """
    pics = []
    named = []
    others = []

    for ph in slide.placeholders:
        try:
            pht = ph.placeholder_format.type
        except Exception:
            pht = None
        name = (getattr(ph, "name", "") or "").lower()

        if pht == PP_PLACEHOLDER.PICTURE:
            pics.append(ph)
        elif "picture" in name or "image" in name:
            named.append(ph)
        else:
            # skip Title/Body for safety
            if pht in (PP_PLACEHOLDER.TITLE, getattr(PP_PLACEHOLDER, "CENTER_TITLE", None), PP_PLACEHOLDER.BODY):
                continue
            # some templates have content placeholders that can hold pictures
            others.append(ph)

    # keep order stable within each group
    return pics + named + others

def right_column_fallback(slide, src_path: Path):
    """Place an image in a right column with safe margins."""
    page_w = slide.part.presentation.slide_width
    page_h = slide.part.presentation.slide_height
    margin = Inches(0.6)
    title_h = Inches(0.9)

    left = int(page_w * 0.58)
    top = margin + title_h
    width = int(page_w - left - margin)
    height = int(page_h - top - margin)

    pic = slide.shapes.add_picture(str(src_path), left, top)
    # scale down to fit
    if pic.width > width or pic.height > height:
        scale = min(width / pic.width, height / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
    # center within that right column
    pic.left = left + int((width - pic.width) / 2)

def apply_images(slide, spec_dir: Path, images):
    """Insert images using picture-like placeholders; leftover images use a right-column fallback."""
    if not images:
        return

    # Resolve all existing image files (preserve order)
    files = []
    for im in images:
        rel = im.get("path")
        if not rel:
            continue
        p = (spec_dir / rel).resolve()
        if p.exists():
            files.append(p)
    if not files:
        return

    # Find target placeholders on this slide
    targets = find_image_placeholders(slide)

    with tempfile.TemporaryDirectory() as td:
        tmpdir = Path(td)

        # 1) Fill available placeholders
        for img_path, ph in zip(files, targets):
            try:
                final_img = ensure_supported_image(img_path, tmpdir)
            except Exception as e:
                print(f"[warn] convert image failed {img_path}: {e}")
                continue

            try:
                # If it's a true PICTURE placeholder, use insert_picture()
                if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE and ph.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    ph.insert_picture(str(final_img))
                else:
                    # Use the placeholder's rect: remove it, then add picture at same location
                    left, top, width, height = ph.left, ph.top, ph.width, ph.height
                    try:
                        slide.shapes._spTree.remove(ph._element)
                    except Exception:
                        pass
                    pic = slide.shapes.add_picture(str(final_img), left, top, width=width)
                    # If aspect mismatch, width is honored and height adjusts. Good enough for template slots.
            except Exception as e:
                print(f"[warn] could not place image in placeholder: {e}")
                right_column_fallback(slide, final_img)

        # 2) Any leftover images → right column fallback (stacked)
        if len(files) > len(targets):
            for img_path in files[len(targets):]:
                try:
                    final_img = ensure_supported_image(img_path, tmpdir)
                except Exception as e:
                    print(f"[warn] convert image failed {img_path}: {e}")
                    continue
                right_column_fallback(slide, final_img)

# ---------------- Build ----------------

def build_from_specs(prs: Presentation, spec_paths, out_file: Path):
    layout = find_preferred_layout(prs)
    count = 0

    for spec_path in spec_paths:
        spec_path = Path(spec_path).resolve()
        spec_dir = spec_path.parent

        for s in iter_slides(spec_path):
            slide = prs.slides.add_slide(layout)
            apply_title(slide, s.get("title","") or "")
            apply_bullets(slide, s.get("bullets", []) or [])
            apply_images(slide, spec_dir, s.get("images", []) or [])

            # Notes (e.g., source URL)
            notes = s.get("notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

            count += 1

    prs.save(str(out_file))
    print(f"[ok] wrote {out_file} with {count} slides (template placeholders)")

# ---------------- CLI ----------------

def main():
    ap = argparse.ArgumentParser(description="Build PPT from slidespec.json using template placeholders.")
    ap.add_argument("--index", help="Path to slidespec_index.csv (preferred)")
    ap.add_argument("--specs", nargs="*", help="Paths to slidespec.json (alternative to --index)")
    ap.add_argument("--template", required=True, help="Path to your PPT/POTX template")
    ap.add_argument("--out", required=True, help="Output PPTX file")
    args = ap.parse_args()

    # Collect slide specs
    if args.index:
        idx = Path(args.index)
        rows = list(csv.DictReader(idx.read_text(encoding="utf-8").splitlines()))
        spec_paths = [Path(r["slidespec"]) for r in rows if r.get("slidespec")]
    elif args.specs:
        spec_paths = [Path(p) for p in args.specs]
    else:
        raise SystemExit("Provide either --index slidespec_index.csv or --specs paths...")

    prs = Presentation(args.template)
    build_from_specs(prs, spec_paths, Path(args.out).resolve())

if __name__ == "__main__":
    main()
